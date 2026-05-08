"""
测试文档处理器

测试覆盖：
- DocumentProcessor 接口
- PPTProcessor 实现
- 文本提取和图片提取
"""

import pytest
from pathlib import Path
from unittest.mock import Mock, MagicMock, patch
from vidppt.core.interfaces import DocumentProcessor
from vidppt.core.models import DocumentContent, PageContent, ProcessConfig


class TestDocumentProcessorInterface:
    """测试 DocumentProcessor 接口"""

    def test_cannot_instantiate_abstract_processor(self):
        """测试不能直接实例化抽象处理器"""
        with pytest.raises(TypeError):
            DocumentProcessor()

    def test_must_implement_all_methods(self):
        """测试必须实现所有抽象方法"""

        class IncompleteProcessor(DocumentProcessor):
            @classmethod
            def supported_extensions(cls):
                return [".test"]

        with pytest.raises(TypeError):
            IncompleteProcessor()

    def test_can_inherit_and_implement(self):
        """测试可以继承并实现接口"""

        class CustomProcessor(DocumentProcessor):
            @classmethod
            def supported_extensions(cls):
                return [".custom"]

            def extract_content(self, config):
                return DocumentContent(pages=[])

            def render_pages(self, config):
                return []

        processor = CustomProcessor()
        assert isinstance(processor, DocumentProcessor)

    def test_process_template_method(self, temp_dir):
        """测试 process 模板方法"""

        class TestProcessor(DocumentProcessor):
            @classmethod
            def supported_extensions(cls):
                return [".test"]

            def extract_content(self, config):
                pages = [
                    PageContent(page_number=1, text="页面1"),
                    PageContent(page_number=2, text="页面2"),
                ]
                return DocumentContent(pages=pages)

            def render_pages(self, config):
                return [
                    Path("slide1.png"),
                    Path("slide2.png"),
                ]

        config = ProcessConfig(
            input_path=temp_dir / "test.test", output_dir=temp_dir / "output"
        )

        processor = TestProcessor()
        content = processor.process(config)

        # 验证模板方法正确组合了提取和渲染
        assert content.total_pages == 2
        assert content.pages[0].slide_image == Path("slide1.png")
        assert content.pages[1].slide_image == Path("slide2.png")


class TestPPTProcessor:
    """测试 PPTProcessor"""

    def test_supported_extensions(self):
        """测试支持的文件扩展名"""
        from vidppt.processors.ppt_processor import PPTProcessor

        extensions = PPTProcessor.supported_extensions()
        assert ".ppt" in extensions
        assert ".pptx" in extensions

    def test_extract_text_from_slide(self):
        """测试从幻灯片提取文本"""
        from vidppt.processors.ppt_processor import PPTProcessor

        # Mock slide with text
        mock_para1 = Mock()
        mock_para1.text = "标题文本"
        mock_para1.level = 0

        mock_para2 = Mock()
        mock_para2.text = "  子标题  "
        mock_para2.level = 1

        mock_shape = Mock()
        mock_shape.has_text_frame = True
        mock_shape.text_frame.paragraphs = [mock_para1, mock_para2]

        mock_slide = Mock()
        mock_slide.shapes = [mock_shape]

        text = PPTProcessor._extract_text_from_slide(mock_slide)

        assert "标题文本" in text
        assert "子标题" in text
        # 验证缩进
        assert "  子标题" in text

    def test_extract_text_skips_empty_paragraphs(self):
        """测试跳过空段落"""
        from vidppt.processors.ppt_processor import PPTProcessor

        mock_para_empty = Mock()
        mock_para_empty.text = "   "
        mock_para_empty.level = 0

        mock_para_valid = Mock()
        mock_para_valid.text = "有效文本"
        mock_para_valid.level = 0

        mock_shape = Mock()
        mock_shape.has_text_frame = True
        mock_shape.text_frame.paragraphs = [mock_para_empty, mock_para_valid]

        mock_slide = Mock()
        mock_slide.shapes = [mock_shape]

        text = PPTProcessor._extract_text_from_slide(mock_slide)

        assert text == "有效文本"

    def test_extract_text_preserves_hierarchy(self):
        """测试保留文本层级"""
        from vidppt.processors.ppt_processor import PPTProcessor

        paragraphs_data = [
            ("一级标题", 0),
            ("二级内容", 1),
            ("三级详情", 2),
            ("另一个一级", 0),
        ]

        mock_shape = Mock()
        mock_shape.has_text_frame = True
        mock_shape.text_frame.paragraphs = []

        for text, level in paragraphs_data:
            para = Mock()
            para.text = text
            para.level = level
            mock_shape.text_frame.paragraphs.append(para)

        mock_slide = Mock()
        mock_slide.shapes = [mock_shape]

        result = PPTProcessor._extract_text_from_slide(mock_slide)
        lines = result.split("\n")

        # 验证缩进
        assert lines[0] == "一级标题"
        assert lines[1] == "  二级内容"
        assert lines[2] == "    三级详情"
        assert lines[3] == "另一个一级"

    def test_extract_images_from_slide(self, temp_dir):
        """测试从幻灯片提取图片"""
        from vidppt.processors.ppt_processor import PPTProcessor
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Mock picture shape
        mock_image = Mock()
        mock_image.ext = "png"
        mock_image.blob = b"fake_image_data"

        mock_picture = Mock()
        mock_picture.shape_type = MSO_SHAPE_TYPE.PICTURE
        mock_picture.image = mock_image

        # Mock other shape
        mock_text_shape = Mock()
        mock_text_shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX

        mock_slide = Mock()
        mock_slide.shapes = [mock_text_shape, mock_picture]

        out_dir = temp_dir / "images"
        out_dir.mkdir(parents=True, exist_ok=True)

        images = PPTProcessor._extract_images_from_slide(mock_slide, out_dir)

        assert len(images) == 1
        assert images[0].exists()
        assert images[0].name == "image_1.png"
        assert images[0].read_bytes() == b"fake_image_data"

    def test_extract_images_from_group(self, temp_dir):
        """测试从组合形状中提取图片"""
        from vidppt.processors.ppt_processor import PPTProcessor
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Mock images in group
        mock_image1 = Mock()
        mock_image1.ext = "jpg"
        mock_image1.blob = b"image1_data"

        mock_picture1 = Mock()
        mock_picture1.shape_type = MSO_SHAPE_TYPE.PICTURE
        mock_picture1.image = mock_image1

        mock_image2 = Mock()
        mock_image2.ext = "png"
        mock_image2.blob = b"image2_data"

        mock_picture2 = Mock()
        mock_picture2.shape_type = MSO_SHAPE_TYPE.PICTURE
        mock_picture2.image = mock_image2

        # Mock group shape
        mock_group = Mock()
        mock_group.shape_type = MSO_SHAPE_TYPE.GROUP
        mock_group.shapes = [mock_picture1, mock_picture2]

        mock_slide = Mock()
        mock_slide.shapes = [mock_group]

        out_dir = temp_dir / "images"
        out_dir.mkdir(parents=True, exist_ok=True)

        images = PPTProcessor._extract_images_from_slide(mock_slide, out_dir)

        assert len(images) == 2
        assert images[0].name == "image_1.jpg"
        assert images[1].name == "image_2.png"


class TestPPTProcessorRenderPages:
    """测试 PPTProcessor render_pages 方法"""

    def test_render_pages_uses_spire_by_default(self, temp_dir):
        """验证默认走 spire 分支"""
        from vidppt.processors.ppt_processor import PPTProcessor

        config = ProcessConfig(
            input_path=temp_dir / "test.pptx",
            output_dir=temp_dir / "output",
            save_intermediate=True,
            skip_existing=False,
        )

        processor = PPTProcessor()

        with patch.object(processor, "_render_with_spire", return_value=[Path("1/slide.png")]) as mock_spire:
            with patch.object(processor, "_render_with_libreoffice") as mock_lo:
                result = processor.render_pages(config)
                mock_spire.assert_called_once_with(config)
                mock_lo.assert_not_called()
                assert result == [Path("1/slide.png")]

    def test_render_pages_uses_libreoffice_when_configured(self, temp_dir):
        """验证 render_engine=libreoffice 时走 libreoffice 分支"""
        from vidppt.processors.ppt_processor import PPTProcessor

        config = ProcessConfig(
            input_path=temp_dir / "test.pptx",
            output_dir=temp_dir / "output",
            render_engine="libreoffice",
            save_intermediate=True,
            skip_existing=False,
        )

        processor = PPTProcessor()

        with patch.object(processor, "_render_with_libreoffice", return_value=[Path("1/slide.png")]) as mock_lo:
            with patch.object(processor, "_render_with_spire") as mock_spire:
                result = processor.render_pages(config)
                mock_lo.assert_called_once_with(config)
                mock_spire.assert_not_called()
                assert result == [Path("1/slide.png")]

    def test_render_with_libreoffice_skip_existing(self, temp_dir):
        """预创建 1/slide.png，验证 skip_existing 生效时不重复渲染该页"""
        from vidppt.processors.ppt_processor import PPTProcessor

        # 创建输入文件
        input_file = temp_dir / "test.pptx"
        input_file.write_text("test")

        # 预创建第1页的 slide.png
        output_dir = temp_dir / "output"
        page1_dir = output_dir / "1"
        page1_dir.mkdir(parents=True, exist_ok=True)
        (page1_dir / "slide.png").write_bytes(b"existing")

        config = ProcessConfig(
            input_path=input_file,
            output_dir=output_dir,
            render_engine="libreoffice",
            save_intermediate=True,
            skip_existing=True,
        )

        processor = PPTProcessor()

        # Mock _get_slide_count 返回 2 页
        with patch.object(processor, "_get_slide_count", return_value=2):
            with patch("vidppt.processors.ppt_processor.subprocess.run") as mock_run:
                def fake_run(cmd, **kwargs):
                    if "pdf" in cmd:
                        # Step 1: LO → PDF (output is always input.pdf now)
                        outdir = cmd[5]
                        (Path(outdir) / "input.pdf").write_bytes(b"fake_pdf")
                    else:
                        # Step 2: pdftoppm → PNG
                        prefix = cmd[-1]  # slide
                        (Path(f"{prefix}-2.png")).write_bytes(b"page2")
                    return MagicMock(returncode=0)

                mock_run.side_effect = fake_run

                result = processor._render_with_libreoffice(config)

                # 应该调用了 subprocess.run（因为第2页不存在）
                assert mock_run.called

    def test_render_with_libreoffice_all_exist_skip(self, temp_dir):
        """所有页的 slide.png 都已存在，验证完全不调用 subprocess.run"""
        from vidppt.processors.ppt_processor import PPTProcessor

        # 创建输入文件
        input_file = temp_dir / "test.pptx"
        input_file.write_text("test")

        # 预创建所有页的 slide.png
        output_dir = temp_dir / "output"
        for i in range(1, 3):
            page_dir = output_dir / str(i)
            page_dir.mkdir(parents=True, exist_ok=True)
            (page_dir / "slide.png").write_bytes(b"existing")

        config = ProcessConfig(
            input_path=input_file,
            output_dir=output_dir,
            render_engine="libreoffice",
            save_intermediate=True,
            skip_existing=True,
        )

        processor = PPTProcessor()

        with patch.object(processor, "_get_slide_count", return_value=2):
            with patch("vidppt.processors.ppt_processor.subprocess.run") as mock_run:
                result = processor._render_with_libreoffice(config)
                mock_run.assert_not_called()
                assert len(result) == 2
                assert all(p.exists() for p in result)

    def test_render_with_libreoffice_file_naming(self, temp_dir):
        """验证 pdftoppm 输出 slide-1.png slide-2.png 被 rename 到 1/slide.png 2/slide.png"""
        from vidppt.processors.ppt_processor import PPTProcessor

        input_file = temp_dir / "演示文稿.pptx"
        input_file.write_text("test")
        output_dir = temp_dir / "output"

        config = ProcessConfig(
            input_path=input_file,
            output_dir=output_dir,
            render_engine="libreoffice",
            save_intermediate=True,
            skip_existing=False,
        )

        processor = PPTProcessor()

        with patch("vidppt.processors.ppt_processor.subprocess.run") as mock_run:
            def fake_run(cmd, **kwargs):
                if "pdf" in cmd:
                    # Step 1: LO → PDF (output is always input.pdf now)
                    outdir = cmd[5]
                    (Path(outdir) / "input.pdf").write_bytes(b"fake_pdf")
                else:
                    # Step 2: pdftoppm → PNG
                    prefix = cmd[-1]  # e.g. /tmp/xxx/slide
                    (Path(f"{prefix}-1.png")).write_bytes(b"page1")
                    (Path(f"{prefix}-2.png")).write_bytes(b"page2")
                return MagicMock(returncode=0)

            mock_run.side_effect = fake_run

            result = processor._render_with_libreoffice(config)

            assert len(result) == 2
            # 验证最终路径
            assert result[0] == output_dir / "1" / "slide.png"
            assert result[1] == output_dir / "2" / "slide.png"
            # 验证文件内容
            assert result[0].read_bytes() == b"page1"
            assert result[1].read_bytes() == b"page2"

    def test_render_engine_not_installed_raises(self, temp_dir):
        """mock subprocess.run 抛 FileNotFoundError，验证错误被正确抛出"""
        from vidppt.processors.ppt_processor import PPTProcessor

        input_file = temp_dir / "test.pptx"
        input_file.write_text("test")
        output_dir = temp_dir / "output"

        config = ProcessConfig(
            input_path=input_file,
            output_dir=output_dir,
            render_engine="libreoffice",
            save_intermediate=True,
            skip_existing=False,
        )

        processor = PPTProcessor()

        with patch("vidppt.processors.ppt_processor.subprocess.run", side_effect=FileNotFoundError("libreoffice not found")):
            with pytest.raises(FileNotFoundError, match="libreoffice not found"):
                processor._render_with_libreoffice(config)


class TestPPTProcessorIntegration:
    """测试 PPTProcessor 集成场景"""

    def test_extract_content_with_mocked_presentation(self, temp_dir):
        """测试使用模拟的演示文稿提取内容"""
        from vidppt.processors.ppt_processor import PPTProcessor

        # Mock slides
        mock_para = Mock()
        mock_para.text = "测试文本"
        mock_para.level = 0

        mock_shape = Mock()
        mock_shape.has_text_frame = True
        mock_shape.text_frame.paragraphs = [mock_para]

        mock_slide = Mock()
        mock_slide.shapes = [mock_shape]

        mock_presentation = Mock()
        mock_presentation.slides = [mock_slide, mock_slide]  # 2个相同的幻灯片

        config = ProcessConfig(
            input_path=temp_dir / "test.pptx",
            output_dir=temp_dir / "output",
            save_intermediate=False,
        )

        processor = PPTProcessor()

        with patch("vidppt.processors.ppt_processor.Presentation") as mock_prs_class:
            mock_prs_class.return_value = mock_presentation

            content = processor.extract_content(config)

            assert content.total_pages == 2
            assert all(p.text == "测试文本" for p in content.pages)

    def test_extract_content_saves_intermediate_files(self, temp_dir):
        """测试提取内容时保存中间文件"""
        from vidppt.processors.ppt_processor import PPTProcessor

        mock_para = Mock()
        mock_para.text = "第一页内容"
        mock_para.level = 0

        mock_shape = Mock()
        mock_shape.has_text_frame = True
        mock_shape.text_frame.paragraphs = [mock_para]

        mock_slide = Mock()
        mock_slide.shapes = [mock_shape]

        mock_presentation = Mock()
        mock_presentation.slides = [mock_slide]

        config = ProcessConfig(
            input_path=temp_dir / "test.pptx",
            output_dir=temp_dir / "output",
            save_intermediate=True,
        )

        processor = PPTProcessor()

        with patch("vidppt.processors.ppt_processor.Presentation") as mock_prs_class:
            mock_prs_class.return_value = mock_presentation

            content = processor.extract_content(config)

            # 验证中间文件已创建
            text_file = temp_dir / "output" / "1" / "text.txt"
            assert text_file.exists()
            assert text_file.read_text(encoding="utf-8") == "第一页内容"

    def test_process_workflow(self, temp_dir):
        """测试完整的处理工作流"""
        from vidppt.processors.ppt_processor import PPTProcessor

        # 这个测试需要模拟整个流程，但由于依赖较多，这里只做基本验证
        processor = PPTProcessor()

        # 验证处理器已正确配置
        assert ".pptx" in processor.supported_extensions()

        # 验证处理器有正确的方法
        assert hasattr(processor, "extract_content")
        assert hasattr(processor, "render_pages")
        assert hasattr(processor, "process")
