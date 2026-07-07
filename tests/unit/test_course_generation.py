"""路线 A：教案到课程模型、字幕和持久化结构。"""

import json
from unittest.mock import patch

import pytest

from vidppt.core.course import Course
from vidppt.core.models import ProcessConfig
from vidppt.generation.course_builder import CourseBuilder
from vidppt.ingestion.models import SourceDocument, SourceSection
from vidppt.renderers.subtitles import SubtitleRenderer


def test_draft_builder_creates_editable_slides(temp_dir):
    source = SourceDocument(
        title="测试课程",
        source_path=temp_dir / "lesson.docx",
        sections=[
            SourceSection(
                title="第一章",
                paragraphs=["理解核心概念；掌握基本方法；完成课堂练习"],
            )
        ],
    )

    course = CourseBuilder().build(source)

    assert course.source_type == "lesson_plan"
    assert course.sections[0].layout == "cover"
    assert course.sections[1].title == "第一章"
    assert course.sections[1].bullets == [
        "理解核心概念",
        "掌握基本方法",
        "完成课堂练习",
    ]
    assert course.sections[1].script


def test_llm_builder_accepts_json_code_fence(temp_dir):
    class FakeLLM:
        def summarize(self, text, **kwargs):
            return """```json
            {
              "title": "AI 课程",
              "audience": "初学者",
              "learning_objectives": ["理解概念"],
              "design_spec": {
                "theme": "technology",
                "mood": "理性、未来感",
                "density": "medium",
                "visual_language": "网络节点与数据流"
              },
              "slides": [{
                "title": "封面",
                "layout": "cover",
                "bullets": [],
                "script": "欢迎学习。",
                "notes": "开场"
              }]
            }
            ```"""

    source = SourceDocument(
        title="原始标题",
        source_path=temp_dir / "lesson.docx",
        sections=[SourceSection(title="正文", paragraphs=["内容"])],
    )

    course = CourseBuilder(FakeLLM()).build(source)

    assert course.title == "AI 课程"
    assert course.audience == "初学者"
    assert course.sections[0].script == "欢迎学习。"
    assert course.metadata["design_spec"]["theme"] == "technology"
    assert course.metadata["design_spec"]["visual_language"] == "网络节点与数据流"


def test_course_new_fields_round_trip():
    payload = {
        "title": "课程",
        "audience": "教师",
        "learning_objectives": ["会使用"],
        "sections": [
            {
                "id": "slide-1",
                "title": "开始",
                "bullets": ["要点"],
                "layout": "cover",
                "notes": "备注",
                "section_title": "导入",
                "script": "讲稿",
            }
        ],
    }

    restored = Course.from_dict(json.loads(json.dumps(Course.from_dict(payload).to_dict())))

    assert restored.audience == "教师"
    assert restored.learning_objectives == ["会使用"]
    assert restored.sections[0].bullets == ["要点"]
    assert restored.sections[0].notes == "备注"


def test_subtitle_renderer_builds_course_timeline(temp_dir):
    course = Course.from_dict(
        {
            "title": "课程",
            "sections": [
                {"id": "1", "title": "一", "script": "第一句。第二句。"},
                {"id": "2", "title": "二", "script": "第三句。"},
            ],
        }
    )

    output = SubtitleRenderer().render_course(
        [(course.sections[0], 4.0), (course.sections[1], 2.0)],
        temp_dir / "course.srt",
    )
    text = output.read_text(encoding="utf-8")

    assert "00:00:00,000" in text
    assert "00:00:04,000 --> 00:00:06,000" in text
    assert "第三句。" in text


def test_subtitle_renderer_splits_long_script_into_single_line_cues(temp_dir):
    course = Course.from_dict(
        {
            "title": "课程",
            "sections": [
                {
                    "id": "1",
                    "title": "一",
                    "script": (
                        "第一，厘清模型定位——企业不是要造大脑，而是选对工具脑；"
                        "第二，直面知识库建设中最容易被忽视的管理难题；"
                        "第三，说明为什么连接CRM、ERP等系统才是真实价值。"
                    ),
                },
            ],
        }
    )

    output = SubtitleRenderer().render_course(
        [(course.sections[0], 12.0)], temp_dir / "course.srt"
    )
    blocks = output.read_text(encoding="utf-8").strip().split("\n\n")

    assert len(blocks) >= 3
    assert all(len(block.splitlines()[2:]) == 1 for block in blocks)
    assert all(
        len(line) <= 20
        for block in blocks
        for line in block.splitlines()[2:]
    )


def test_course_pipeline_rejects_video_without_tts(temp_dir):
    from vidppt.core.models import ProcessConfig
    from vidppt.course_pipeline import CoursePipeline

    source = temp_dir / "lesson.docx"
    source.write_bytes(b"placeholder")
    config = ProcessConfig(
        input_path=source,
        output_dir=temp_dir / "out",
        enable_tts=False,
        enable_video=True,
    )

    with pytest.raises(ValueError, match="不能关闭 TTS"):
        CoursePipeline().run(source, config.output_dir, media_config=config)


def test_course_pipeline_burns_subtitles_into_video(temp_dir):
    from vidppt.course_pipeline import CoursePipeline

    video = temp_dir / "base.mp4"
    subtitles = temp_dir / "course.srt"
    output = temp_dir / "course.mp4"

    with (
        patch("vidppt.course_pipeline.shutil.which", return_value="/usr/bin/ffmpeg"),
        patch("vidppt.course_pipeline.subprocess.run") as run,
    ):
        CoursePipeline._burn_subtitles(
            video,
            subtitles,
            output,
            ProcessConfig(input_path=video, output_dir=temp_dir),
        )

    command = run.call_args.args[0]
    assert "-vf" in command
    assert "subtitles=filename='course.srt'" in command[command.index("-vf") + 1]
    subtitle_filter = command[command.index("-vf") + 1]
    assert "drawbox=x=0:y=976:w=1920:h=50" in subtitle_filter
    assert "color=0x333333@0.45" in subtitle_filter
    assert "FontSize=50" in subtitle_filter
    assert "PrimaryColour=&H00FFFFFF" in subtitle_filter
    assert "Outline=0" in subtitle_filter
    assert "Alignment=2" in subtitle_filter
    assert command[command.index("-c:v") + 1] == "libx264"
    assert command[command.index("-preset") + 1] == "veryfast"
    assert command[command.index("-crf") + 1] == "21"
    assert command[command.index("-pix_fmt") + 1] == "yuv420p"
    assert command[command.index("-g") + 1] == "48"
    assert command[command.index("-movflags") + 1] == "+faststart"
    assert run.call_args.kwargs["cwd"] == temp_dir


def test_moviepy_progress_uses_video_frames_not_audio_chunks():
    from vidppt.utils.video_composer import _MoviePyProgressLogger

    reported = []
    progress = _MoviePyProgressLogger.create(reported.append, None)
    progress.bars["chunk"] = {"total": 10}
    progress.bars_callback("chunk", "index", 10)
    assert reported == []

    progress.bars["frame_index"] = {"total": 200}
    progress.bars_callback("frame_index", "index", 50)
    assert reported == [0.25]


def test_docx_pipeline_outputs_course_json_and_editable_pptx(temp_dir):
    from docx import Document
    from pptx import Presentation

    from vidppt.course_pipeline import CoursePipeline

    source = temp_dir / "lesson.docx"
    document = Document()
    document.add_heading("Shell 脚本入门", 0)
    document.add_heading("变量", 1)
    document.add_paragraph("理解变量定义；掌握变量引用")
    document.save(source)

    result = CoursePipeline(footer_text="匠心课堂").run(source, temp_dir / "out")

    assert result.course_json.exists()
    assert result.presentation.exists()
    payload = json.loads(result.course_json.read_text(encoding="utf-8"))
    assert payload["title"] == "Shell 脚本入门"
    assert payload["metadata"]["footer_text"] == "匠心课堂"
    assert payload["sections"][1]["bullets"] == ["理解变量定义", "掌握变量引用"]

    presentation = Presentation(result.presentation)
    assert len(presentation.slides) == 2
    slide_text = "\n".join(
        shape.text for shape in presentation.slides[1].shapes if hasattr(shape, "text")
    )
    assert "变量" in slide_text
    assert "理解变量定义" in slide_text
    assert "匠心课堂" in slide_text
    assert "理解变量定义" in presentation.slides[1].notes_slide.notes_text_frame.text


def test_course_pipeline_places_logo_without_cropping(temp_dir):
    from docx import Document
    from PIL import Image
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from vidppt.course_pipeline import CoursePipeline

    source = temp_dir / "lesson.docx"
    document = Document()
    document.add_heading("学校课程", 0)
    document.add_heading("课程内容", 1)
    document.add_paragraph("知识点一；知识点二")
    document.save(source)
    logo = temp_dir / "school.png"
    Image.new("RGBA", (400, 100), (20, 80, 160, 255)).save(logo)

    result = CoursePipeline(logo_path=logo).run(source, temp_dir / "out-logo")
    payload = json.loads(result.course_json.read_text(encoding="utf-8"))
    assert payload["metadata"]["logo_path"] == str(logo)

    presentation = Presentation(result.presentation)
    for slide in presentation.slides:
        pictures = [
            shape for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        assert pictures
        logo_shape = pictures[-1]
        assert abs((logo_shape.width / logo_shape.height) - 4.0) < 0.05
        assert logo_shape.width <= 1.35 * 914400
        assert logo_shape.height <= 0.72 * 914400


def test_course_pipeline_generates_illustrations_before_pptx(temp_dir):
    from docx import Document
    from vidppt.course_pipeline import CoursePipeline

    source = temp_dir / "lesson.docx"
    document = Document()
    document.add_heading("智能制造", 0)
    document.add_heading("精密加工", 1)
    document.add_paragraph("理解加工精度；掌握质量控制")
    document.save(source)

    class FakeIllustrations:
        def generate_for_course(self, course, output_dir, max_images):
            assert max_images == 2
            course.sections[0].metadata["illustration_path"] = "cover.png"
            return []

    result = CoursePipeline(
        illustration_generator=FakeIllustrations(),
        max_illustrations=2,
    ).run(source, temp_dir / "out")
    payload = json.loads(result.course_json.read_text(encoding="utf-8"))

    assert payload["metadata"]["visual_theme"] == "industry"
    assert payload["metadata"]["visual_style"] == "technical"
    assert payload["sections"][0]["metadata"]["illustration_path"] == "cover.png"


def test_pptx_theme_selects_distinct_visual_systems():
    from vidppt.core.course import Course, CourseSection
    from vidppt.renderers.pptx_renderer import PPTXRenderer

    samples = {
        "technology": ("人工智能与数据系统", "digital"),
        "culture": ("中国传统文学艺术", "editorial"),
        "nature": ("生态环境与自然保护", "organic"),
        "business": ("企业营销与运营管理", "executive"),
        "health": ("医疗护理与健康管理", "organic"),
        "public": ("公共政策与社会治理", "executive"),
        "finance": ("财务审计与证券投资", "executive"),
    }
    for expected_theme, (title, expected_style) in samples.items():
        course = Course(
            title=title,
            sections=[CourseSection(id="slide-1", title=title, layout="cover")],
        )
        renderer = PPTXRenderer()
        renderer._apply_course_theme(course)
        assert course.metadata["visual_theme"] == expected_theme
        assert course.metadata["visual_style"] == expected_style


def test_manual_theme_overrides_content_and_auto_layouts_vary():
    from vidppt.core.course import Course, CourseSection
    from vidppt.renderers.pptx_renderer import PPTXRenderer

    course = Course(
        title="人工智能系统",
        metadata={"visual_theme": "culture"},
        sections=[
            CourseSection(id="1", title="封面", layout="cover"),
            CourseSection(
                id="2",
                title="系统实施流程",
                bullets=["需求分析", "方案设计", "部署验证"],
            ),
            CourseSection(
                id="3",
                title="架构的四大支柱",
                bullets=["模型", "知识", "系统", "治理"],
            ),
        ],
    )
    renderer = PPTXRenderer()
    renderer._apply_course_theme(course)

    assert course.metadata["visual_theme"] == "culture"
    assert renderer._resolve_layout(course.sections[1], 2) == "process"
    assert renderer._resolve_layout(course.sections[2], 3) == "framework"
