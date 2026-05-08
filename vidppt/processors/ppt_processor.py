"""
PowerPoint 文档处理器
"""

import io
import shutil
import subprocess
import tempfile
from pathlib import Path

from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from ..core.interfaces import DocumentProcessor
from ..core.models import DocumentContent, PageContent, ProcessConfig
from ..core.registry import register_processor


@register_processor
class PPTProcessor(DocumentProcessor):
    """PowerPoint 文档处理器"""

    @classmethod
    def supported_extensions(cls) -> list[str]:
        return [".ppt", ".pptx"]

    def extract_content(self, config: ProcessConfig) -> DocumentContent:
        """提取 PPT 内容"""
        prs = Presentation(str(config.input_path))
        pages = []

        logger.info(f"共 {len(prs.slides)} 页，开始提取...")

        text_skipped = 0
        for i, slide in enumerate(prs.slides, start=1):
            page = PageContent(page_number=i)
            page_dir = config.output_dir / str(i)
            text_path = page_dir / "text.txt"

            # 跳过已有文本的页面
            if config.save_intermediate and config.skip_existing and text_path.exists():
                page.text = text_path.read_text(encoding="utf-8")
                logger.debug(f"第 {i} 页 文字已存在，跳过提取  ({len(page.text)} 字符)")
                text_skipped += 1
            else:
                # 提取文本
                page.text = self._extract_text_from_slide(slide)

                # 保存文本（如果需要且不为空）
                if config.save_intermediate and page.text and page.text.strip():
                    page_dir.mkdir(parents=True, exist_ok=True)
                    text_path.write_text(page.text, encoding="utf-8")
                    logger.debug(f"第 {i} 页 文字 -> {text_path}  ({len(page.text)} 字符)")
                elif config.save_intermediate:
                    logger.debug(f"第 {i} 页 无文本内容，跳过文本文件保存")

                # 提取图片
                if config.save_intermediate:
                    page_dir.mkdir(parents=True, exist_ok=True)
                    page.images = self._extract_images_from_slide(slide, page_dir)
                    for img_path in page.images:
                        logger.debug(f"第 {i} 页 图片 -> {img_path}")
                    if not page.images:
                        logger.debug(f"第 {i} 页 无内嵌图片")

            pages.append(page)

        if text_skipped:
            logger.info(f"跳过 {text_skipped} 页已有文字提取")

        return DocumentContent(pages=pages)

    def render_pages(self, config: ProcessConfig) -> list[Path]:
        """渲染 PPT 页面为图像，根据 render_engine 分发"""
        if config.render_engine == "libreoffice":
            return self._render_with_libreoffice(config)
        return self._render_with_spire(config)

    def _render_with_spire(self, config: ProcessConfig) -> list[Path]:
        """使用 Spire 渲染 PPT 页面为图像"""
        from spire.presentation import Presentation as SpirePresentation

        logger.info("开始渲染幻灯片截图（spire）...")

        # 预检查哪些页需要渲染
        if config.save_intermediate and config.skip_existing:
            page_count = self._get_slide_count(config)
            skip_flags = []
            for i in range(1, page_count + 1):
                slide_path = config.output_dir / str(i) / "slide.png"
                skip_flags.append(slide_path.exists())

            if all(skip_flags):
                logger.info("所有幻灯片截图已存在，跳过渲染")
                return [config.output_dir / str(i) / "slide.png" for i in range(1, page_count + 1)]
        else:
            skip_flags = []

        # 需要渲染至少一页，加载 spire
        prs = SpirePresentation()
        prs.LoadFromFile(str(config.input_path))

        # 确保输出目录存在
        config.output_dir.mkdir(parents=True, exist_ok=True)

        # 如果没有预检查结果（skip_existing 未启用），重新计算
        if len(skip_flags) != prs.Slides.Count:
            skip_flags = [False] * prs.Slides.Count
            if config.save_intermediate and config.skip_existing:
                for i in range(prs.Slides.Count):
                    slide_path = config.output_dir / str(i + 1) / "slide.png"
                    skip_flags[i] = slide_path.exists()

        slide_images = []
        rendered = 0
        for i in range(prs.Slides.Count):
            page_num = i + 1

            if skip_flags[i]:
                out_path = config.output_dir / str(page_num) / "slide.png"
                slide_images.append(out_path)
                logger.debug(f"第 {page_num} 页 截图已存在，跳过渲染")
                continue

            slide = prs.Slides[i]
            img_stream = slide.SaveAsImage()
            img_bytes = bytes(img_stream.ToArray())
            img = Image.open(io.BytesIO(img_bytes))

            if config.save_intermediate:
                page_dir = config.output_dir / str(page_num)
                page_dir.mkdir(parents=True, exist_ok=True)
                out_path = page_dir / "slide.png"
                img.save(out_path, "PNG")
                logger.debug(f"第 {page_num} 页 截图 -> {out_path}")
                slide_images.append(out_path)
                rendered += 1
            else:
                # 不保存中间文件时，保存到临时路径
                config.output_dir.mkdir(parents=True, exist_ok=True)
                temp_path = config.output_dir / f"_temp_slide_{page_num}.png"
                img.save(temp_path, "PNG")
                slide_images.append(temp_path)
                rendered += 1

            img_stream.Dispose()

        prs.Dispose()

        skipped = sum(skip_flags)
        if skipped:
            logger.info(f"跳过 {skipped} 页已有截图，渲染 {rendered} 页")
        else:
            logger.info(f"渲染 {rendered} 页截图")

        return slide_images

    def _render_with_libreoffice(self, config: ProcessConfig) -> list[Path]:
        """使用 LibreOffice 渲染 PPT 页面为图像

        两步流程：LibreOffice 转 PDF → pdftoppm 逐页转 PNG
        """
        logger.info("开始渲染幻灯片截图（libreoffice）...")

        # 预检查 skip_existing
        if config.save_intermediate and config.skip_existing:
            page_count = self._get_slide_count(config)
            skip_flags = []
            for i in range(1, page_count + 1):
                slide_path = config.output_dir / str(i) / "slide.png"
                skip_flags.append(slide_path.exists())

            if all(skip_flags):
                logger.info("所有幻灯片截图已存在，跳过渲染")
                return [config.output_dir / str(i) / "slide.png" for i in range(1, page_count + 1)]

            need_render = not all(skip_flags)
        else:
            need_render = True

        if need_render:
            with tempfile.TemporaryDirectory() as tmp_dir:
                # 先将输入文件复制到临时目录（安全文件名），避免 LO 无法加载
                # 含特殊字符（引号等）的原始文件名或某些 PPTX 格式会导致 LO 加载失败
                safe_input = Path(tmp_dir) / "input.pptx"
                shutil.copy2(config.input_path, safe_input)

                # Step 1: LibreOffice PPT → PDF
                cmd_pdf = [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmp_dir,
                    str(safe_input),
                ]
                try:
                    subprocess.run(cmd_pdf, check=True, capture_output=True)
                except subprocess.CalledProcessError:
                    # 某些 PPTX 格式 LO 无法直接加载，用 python-pptx 重新保存后再试
                    logger.info("LibreOffice 加载原始文件失败，尝试重新保存后转换...")
                    prs = Presentation(str(config.input_path))
                    prs.save(str(safe_input))
                    subprocess.run(cmd_pdf, check=True, capture_output=True)

                pdf_path = Path(tmp_dir) / "input.pdf"
                if not pdf_path.exists():
                    raise FileNotFoundError(
                        f"LibreOffice 未生成 PDF: {pdf_path}"
                    )

                # Step 2: pdftoppm PDF → PNG（逐页）
                cmd_ppm = [
                    "pdftoppm",
                    "-png",
                    "-r",
                    "150",
                    str(pdf_path),
                    str(Path(tmp_dir) / "slide"),
                ]
                subprocess.run(cmd_ppm, check=True, capture_output=True)

                # pdftoppm 输出命名：slide-01.png, slide-02.png, ...（零填充）
                config.output_dir.mkdir(parents=True, exist_ok=True)
                slide_images = []
                page_num = 1
                while True:
                    # pdftoppm 用零填充页码，如 slide-01.png, slide-001.png
                    # 两种格式都检查
                    ppm_output = Path(tmp_dir) / f"slide-{page_num:02d}.png"
                    if not ppm_output.exists():
                        # 也尝试三位零填充
                        ppm_output = Path(tmp_dir) / f"slide-{page_num:03d}.png"
                    if not ppm_output.exists():
                        # 最后尝试无零填充
                        ppm_output = Path(tmp_dir) / f"slide-{page_num}.png"
                    if not ppm_output.exists():
                        break
                    page_dir = config.output_dir / str(page_num)
                    page_dir.mkdir(parents=True, exist_ok=True)
                    dest = page_dir / "slide.png"
                    shutil.copy2(ppm_output, dest)
                    slide_images.append(dest)
                    logger.debug(f"第 {page_num} 页 截图 -> {dest}")
                    page_num += 1
        else:
            # 部分存在，只需收集已存在的
            page_count = self._get_slide_count(config)
            slide_images = []
            for i in range(1, page_count + 1):
                slide_path = config.output_dir / str(i) / "slide.png"
                slide_images.append(slide_path)

        logger.info(f"渲染 {len(slide_images)} 页截图（libreoffice）")
        return slide_images

    def _get_slide_count(self, config: ProcessConfig) -> int:
        """获取 PPT 幻灯片数量（不加载 spire）"""
        prs = Presentation(str(config.input_path))
        count = len(prs.slides)
        return count

    @staticmethod
    def _extract_text_from_slide(slide) -> str:
        """提取幻灯片中所有文本框的文字，保留层级缩进"""
        lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                level = para.level if para.level else 0
                indent = "  " * level
                lines.append(f"{indent}{text}")
        return "\n".join(lines)

    @staticmethod
    def _extract_images_from_slide(slide, out_dir: Path) -> list[Path]:
        """提取幻灯片中所有内嵌图片，返回保存路径列表"""
        saved = []
        img_index = 1

        def extract_from_shape(shape):
            nonlocal img_index
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                img_path = out_dir / f"image_{img_index}.{ext}"
                # 如果目录不存在，先创建
                img_path.parent.mkdir(parents=True, exist_ok=True)
                img_path.write_bytes(image.blob)
                saved.append(img_path)
                img_index += 1

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                extract_from_shape(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        extract_from_shape(child)

        return saved
