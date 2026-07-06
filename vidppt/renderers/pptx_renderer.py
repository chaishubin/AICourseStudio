"""将 Course 渲染为具有课程主题、可编辑文本和逐页备注的 PPTX。"""

from pathlib import Path

from ..core.course import Course, CourseSection


class PPTXRenderer:
    """内置课程视觉主题，避免依赖客户端默认的黑白 PowerPoint 母版。"""

    NAVY = "10233F"
    BLUE = "2457D6"
    TEAL = "13A89E"
    ORANGE = "F5A623"
    PAPER = "F6F8FC"
    WHITE = "FFFFFF"
    INK = "172033"
    MUTED = "60708A"
    LINE = "DCE3EF"
    FONT = "Noto Sans CJK SC"
    THEMES = {
        "industry": ("16212F", "276B8A", "E28B32", "F4F6F7"),
        "technology": ("101D42", "2864DC", "16B8A6", "F4F7FC"),
        "culture": ("3B2034", "A34A5B", "D8A23D", "FBF7F2"),
        "nature": ("173D35", "2E8B70", "D6A84B", "F3F8F4"),
        "education": ("243052", "5967C9", "F1A545", "F7F7FC"),
        "business": ("1D2B3A", "2E708C", "D99B3D", "F5F7F9"),
        "health": ("183B45", "2A8C82", "E98575", "F3F9F8"),
        "public": ("243B53", "486581", "D4A72C", "F5F7FA"),
        "finance": ("1C2B33", "2F6B5F", "C7A24B", "F6F7F4"),
    }
    THEME_STYLES = {
        "industry": "technical",
        "technology": "digital",
        "culture": "editorial",
        "nature": "organic",
        "education": "academic",
        "business": "executive",
        "health": "organic",
        "public": "executive",
        "finance": "executive",
    }

    def render(self, course: Course, output_path: Path) -> Path:
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        presentation.core_properties.title = course.title
        presentation.core_properties.subject = course.description
        self._apply_course_theme(course)
        self.footer_text = str(
            course.metadata.get("footer_text", "AI COURSE STUDIO")
        ).strip()[:40]

        total = len(course.sections)
        for index, page in enumerate(course.sections, 1):
            self._add_slide(presentation, course, page, index, total)

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        return output_path

    def _add_slide(
        self,
        presentation,
        course: Course,
        page: CourseSection,
        index: int,
        total: int,
    ) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        layout = (page.layout or "title_and_content").lower()
        if layout == "cover" or index == 1:
            self._render_cover(slide, course, page)
            logo_layout = "cover"
        elif layout == "section":
            self._render_section(slide, page, index, total)
            logo_layout = "section"
        elif layout == "summary" or index == total:
            self._render_summary(slide, page, index, total)
            logo_layout = "summary"
        else:
            self._render_content(slide, page, index, total)
            logo_layout = "content"
        self._add_logo(slide, course, logo_layout)
        self._write_notes(slide, page)

    def _render_cover(self, slide, course: Course, page: CourseSection) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN

        illustration = self._illustration_path(page)
        style = self.DESIGN_STYLE
        light_cover = style in {"editorial", "organic"}

        self._background(slide, self.PAPER if light_cover else self.NAVY)
        if style == "editorial":
            self._shape(slide, MSO_SHAPE.RECTANGLE, 9.55, 0, 3.78, 7.5, self.NAVY)
            self._shape(slide, MSO_SHAPE.RECTANGLE, 0.85, 0.8, 0.08, 5.8, self.ORANGE)
            self._shape(slide, MSO_SHAPE.RECTANGLE, 9.1, 0.8, 0.03, 5.8, self.LINE)
        elif style == "organic":
            self._shape(slide, MSO_SHAPE.OVAL, 9.25, -0.65, 4.7, 4.7, self.TEAL, 0.1)
            self._shape(slide, MSO_SHAPE.OVAL, 10.2, 4.55, 3.2, 3.2, self.ORANGE, 0.2)
            self._shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0, 0, 0.22, 7.5, self.TEAL)
        elif style == "executive":
            self._shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 4.15, 7.5, self.BLUE)
            self._shape(slide, MSO_SHAPE.RECTANGLE, 4.15, 0, 0.1, 7.5, self.ORANGE)
            self._shape(slide, MSO_SHAPE.RECTANGLE, 10.75, 0.7, 1.8, 0.08, self.ORANGE)
        else:
            self._shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, self.TEAL)
            self._shape(slide, MSO_SHAPE.RECTANGLE, 0.18, 0, 0.07, 7.5, self.ORANGE)
            self._shape(slide, MSO_SHAPE.OVAL, 10.95, 0.0, 2.35, 2.35, self.BLUE, 0.35)
            self._shape(slide, MSO_SHAPE.OVAL, 11.2, 5.4, 2.1, 2.1, self.TEAL, 0.3)

        if illustration:
            if style == "editorial":
                self._picture(slide, illustration, 9.55, 0, 3.78, 7.5)
            elif style == "organic":
                self._picture(slide, illustration, 9.05, 1.0, 3.55, 5.5)
            else:
                self._picture(slide, illustration, 8.55, 0, 4.78, 7.5)

        self._text(
            slide, self.footer_text, 0.9, 0.75, 4.5, 0.35,
            13, self.TEAL if not light_cover else self.BLUE, bold=True, letter_spacing=1.5,
        )
        self._text(
            slide, page.title or course.title, 1.15 if style == "editorial" else 0.9,
            2.05, 7.25 if illustration else (7.7 if style == "editorial" else 10.4), 2.0,
            34, self.INK if light_cover else self.WHITE, bold=True, align=PP_ALIGN.LEFT,
        )
        subtitle = course.description or page.section_title or "结构化课程 · 可编辑课件 · 智能讲稿"
        self._text(
            slide, subtitle, 1.15 if style == "editorial" else 0.95, 4.35, 7.9, 0.7,
            18, self.MUTED if light_cover else "C8D5EA",
        )
        self._shape(slide, MSO_SHAPE.RECTANGLE, 1.15 if style == "editorial" else 0.95, 5.42, 1.0, 0.08, self.ORANGE)
        self._text(
            slide, "课程导入", 1.15 if style == "editorial" else 0.95, 5.68, 2.5, 0.4,
            14, self.INK if light_cover else self.WHITE, bold=True,
        )

    def _render_section(
        self, slide, page: CourseSection, index: int, total: int
    ) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN

        self._background(slide, self.BLUE)
        self._shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.22, 7.5, self.ORANGE)
        self._text(
            slide, f"{index:02d}", 0.9, 1.0, 2.0, 1.2,
            42, "8FB0FF", bold=True,
        )
        self._text(
            slide, page.title, 0.9, 2.25, 10.7, 1.4,
            34, self.WHITE, bold=True, align=PP_ALIGN.LEFT,
        )
        if page.bullets:
            self._text(slide, page.bullets[0], 0.95, 4.05, 9.6, 0.8, 19, "DCE7FF")
        self._footer(slide, index, total, dark=True)

    def _render_content(
        self, slide, page: CourseSection, index: int, total: int
    ) -> None:
        from pptx.enum.shapes import MSO_SHAPE

        self._background(slide, self.PAPER)
        accent = self.TEAL if index % 2 else self.BLUE
        illustration = self._illustration_path(page)
        if self.DESIGN_STYLE == "executive":
            self._shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 2.75, 7.5, self.NAVY)
            self._shape(slide, MSO_SHAPE.RECTANGLE, 2.75, 0, 0.09, 7.5, self.ORANGE)
            self._text(slide, page.section_title or "课程内容", 0.45, 0.8, 1.8, 0.4, 12, self.ORANGE, bold=True)
            self._text(slide, page.title, 0.45, 1.55, 1.9, 3.8, 25, self.WHITE, bold=True)
            if illustration:
                self._picture(slide, illustration, 9.25, 1.0, 3.45, 5.65)
            self._render_content_cards(
                slide, page, 3.25, 0.9, 5.65 if illustration else 9.25,
                5.85, accent, columns=1 if illustration else 2,
            )
            self._footer(slide, index, total)
            return
        if self.DESIGN_STYLE in {"editorial", "organic"}:
            self._shape(
                slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.16,
                self.ORANGE if self.DESIGN_STYLE == "editorial" else self.TEAL,
            )
            self._text(slide, page.section_title or "课程内容", 0.78, 0.6, 2.5, 0.3, 11, accent, bold=True)
            self._text(slide, page.title, 0.78, 1.03, 11.5, 0.8, 29, self.INK, bold=True)
            if illustration:
                self._picture(slide, illustration, 8.85, 2.15, 3.68, 4.45)
            self._render_content_cards(
                slide, page, 0.78, 2.15, 7.72 if illustration else 11.75,
                4.45, accent,
                columns=1 if illustration else (2 if self.DESIGN_STYLE == "editorial" else 3),
            )
            self._footer(slide, index, total)
            return

        self._shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, accent)
        self._shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.75, 0.55, 1.7, 0.42, accent)
        self._text(
            slide, page.section_title or "课程内容", 10.88, 0.61, 1.4, 0.24,
            10, self.WHITE, bold=True,
        )
        self._text(slide, page.title, 0.75, 0.72, 9.7, 0.82, 27, self.INK, bold=True)
        self._shape(slide, MSO_SHAPE.RECTANGLE, 0.78, 1.72, 0.85, 0.07, self.ORANGE)

        if illustration:
            self._picture(slide, illustration, 8.65, 2.05, 3.85, 3.9)
            self._shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.65, 6.05, 3.85, 0.52, self.WHITE)
            self._text(slide, "课程概念插图", 8.9, 6.18, 3.3, 0.2, 10, self.MUTED, bold=True)

        bullets = page.bullets[:5 if illustration else 6] or [page.title]
        bullet_width = 6.8 if illustration else 10.6
        top = 2.05
        for bullet_index, bullet in enumerate(bullets, 1):
            y = top + (bullet_index - 1) * 0.78
            self._shape(slide, MSO_SHAPE.OVAL, 0.82, y + 0.08, 0.35, 0.35, accent)
            self._text(
                slide, str(bullet_index), 0.82, y + 0.1, 0.35, 0.24,
                10, self.WHITE, bold=True, align=None,
            )
            self._text(slide, bullet, 1.38, y, bullet_width, 0.58, 20, self.INK)

        self._footer(slide, index, total)

    def _render_content_cards(
        self, slide, page: CourseSection, x: float, y: float,
        width: float, height: float, accent: str, columns: int,
    ) -> None:
        """按主题将正文组织为编辑式、自然式或商务式信息卡。"""
        from pptx.enum.shapes import MSO_SHAPE

        bullets = page.bullets[:6] or [page.title]
        rows = (len(bullets) + columns - 1) // columns
        gap_x, gap_y = 0.22, 0.22
        card_width = (width - gap_x * (columns - 1)) / columns
        card_height = min(1.7, (height - gap_y * max(0, rows - 1)) / max(1, rows))
        for bullet_index, bullet in enumerate(bullets):
            col, row = bullet_index % columns, bullet_index // columns
            card_x = x + col * (card_width + gap_x)
            card_y = y + row * (card_height + gap_y)
            self._shape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                card_x, card_y, card_width, card_height, self.WHITE,
            )
            marker = self.ORANGE if bullet_index % 2 else accent
            if self.DESIGN_STYLE == "organic":
                self._shape(slide, MSO_SHAPE.OVAL, card_x + 0.25, card_y + 0.28, 0.42, 0.42, marker)
            else:
                self._shape(slide, MSO_SHAPE.RECTANGLE, card_x, card_y, 0.08, card_height, marker)
            self._text(
                slide, f"{bullet_index + 1:02d}",
                card_x + 0.25, card_y + 0.22, 0.48, 0.35,
                11, marker, bold=True,
            )
            self._text(
                slide, bullet, card_x + 0.78, card_y + 0.2,
                card_width - 1.0, card_height - 0.38,
                16 if columns == 3 else 18, self.INK, bold=columns == 2,
            )

    def _render_summary(
        self, slide, page: CourseSection, index: int, total: int
    ) -> None:
        from pptx.enum.shapes import MSO_SHAPE

        self._background(slide, self.PAPER)
        self._shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 1.55, self.NAVY)
        self._text(slide, page.title or "课程小结", 0.8, 0.48, 10.8, 0.65, 29, self.WHITE, bold=True)
        bullets = page.bullets[:4] or ["回顾本节核心内容"]
        card_width = 5.7
        for bullet_index, bullet in enumerate(bullets):
            col, row = bullet_index % 2, bullet_index // 2
            x, y = 0.75 + col * 6.15, 2.0 + row * 1.75
            accent = self.TEAL if bullet_index % 2 == 0 else self.ORANGE
            self._shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, 1.35, self.WHITE)
            self._shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.09, 1.35, accent)
            self._text(slide, bullet, x + 0.35, y + 0.28, 5.0, 0.78, 18, self.INK, bold=True)
        self._footer(slide, index, total)

    def _background(self, slide, color: str) -> None:
        from pptx.dml.color import RGBColor

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color)

    def _apply_course_theme(self, course: Course) -> None:
        text = " ".join(
            [
                course.title,
                course.description,
                *(section.title for section in course.sections),
                *(bullet for section in course.sections for bullet in section.bullets),
            ]
        ).lower()
        keyword_groups = {
            "industry": ("工匠", "制造", "工业", "机械", "工程", "设备", "技能", "生产"),
            "technology": ("ai", "人工智能", "编程", "代码", "数据", "科技", "软件", "网络"),
            "culture": ("文化", "历史", "文学", "艺术", "思想", "传统", "精神"),
            "nature": ("生态", "环境", "自然", "农业", "能源", "地理", "生物"),
            "business": ("商业", "管理", "营销", "市场", "企业", "战略", "运营"),
            "health": ("医疗", "健康", "护理", "医学", "疾病", "药物", "营养", "心理"),
            "public": ("党建", "政策", "法律", "政务", "公共", "安全", "治理", "法规"),
            "finance": ("财务", "会计", "投资", "证券", "银行", "审计", "税务", "资本"),
        }
        scores = {
            theme: sum(text.count(keyword) for keyword in keywords)
            for theme, keywords in keyword_groups.items()
        }
        theme_name = max(scores, key=scores.get) if any(scores.values()) else "education"
        navy, blue, accent, paper = self.THEMES[theme_name]
        self.NAVY, self.BLUE, self.TEAL, self.PAPER = navy, blue, accent, paper
        self.DESIGN_STYLE = self.THEME_STYLES[theme_name]
        course.metadata["visual_theme"] = theme_name
        course.metadata["visual_style"] = self.DESIGN_STYLE

    @staticmethod
    def _illustration_path(page: CourseSection) -> Path | None:
        value = page.metadata.get("illustration_path")
        path = Path(value) if value else None
        return path if path and path.exists() else None

    @staticmethod
    def _picture(slide, path: Path, x, y, width, height):
        from PIL import Image
        from pptx.util import Inches

        with Image.open(path) as image:
            image_ratio = image.width / image.height
        frame_ratio = width / height
        picture = slide.shapes.add_picture(
            str(path), Inches(x), Inches(y), Inches(width), Inches(height)
        )
        if image_ratio > frame_ratio:
            visible = frame_ratio / image_ratio
            picture.crop_left = picture.crop_right = (1 - visible) / 2
        elif image_ratio < frame_ratio:
            visible = image_ratio / frame_ratio
            picture.crop_top = picture.crop_bottom = (1 - visible) / 2
        return picture

    def _shape(
        self, slide, shape_type, x, y, width, height, color: str, transparency=0
    ):
        from pptx.dml.color import RGBColor
        from pptx.util import Inches

        shape = slide.shapes.add_shape(
            shape_type, Inches(x), Inches(y), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color)
        shape.fill.transparency = int(transparency * 100)
        shape.line.fill.background()
        return shape

    def _text(
        self,
        slide,
        text,
        x,
        y,
        width,
        height,
        size,
        color,
        *,
        bold=False,
        align=None,
        letter_spacing=None,
    ):
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        box = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.text = str(text)
        paragraph.alignment = align or PP_ALIGN.LEFT
        paragraph.space_after = Pt(0)
        run = paragraph.runs[0]
        run.font.name = self.FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)
        return box

    def _footer(self, slide, index: int, total: int, dark: bool = False) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN

        color = "D7E2F5" if dark else self.MUTED
        line = "6E93E8" if dark else self.LINE
        self._shape(slide, MSO_SHAPE.RECTANGLE, 0.78, 7.05, 11.75, 0.02, line)
        self._text(
            slide, self.footer_text, 0.78, 7.12, 4.8, 0.2,
            9, color, bold=True,
        )
        self._text(
            slide, f"{index:02d} / {total:02d}", 11.2, 7.12, 1.3, 0.2,
            9, color, bold=True, align=PP_ALIGN.RIGHT,
        )

    def _add_logo(self, slide, course: Course, layout: str) -> None:
        """在各版式的安全区内等比放置 Logo，不裁剪且不侵占正文。"""
        from PIL import Image
        from pptx.util import Inches

        value = course.metadata.get("logo_path")
        path = Path(value) if value else None
        if not path or not path.is_file():
            return
        safe_areas = {
            "cover": (11.15, 0.55, 1.35, 0.72),
            "section": (11.15, 0.62, 1.35, 0.72),
            "content": (11.10, 1.08, 1.25, 0.62),
            "summary": (11.20, 0.43, 1.20, 0.66),
        }
        x, y, max_width, max_height = safe_areas[layout]
        with Image.open(path) as image:
            ratio = image.width / image.height
        width = min(max_width, max_height * ratio)
        height = width / ratio
        if height > max_height:
            height = max_height
            width = height * ratio
        x += max_width - width
        y += (max_height - height) / 2
        slide.shapes.add_picture(
            str(path), Inches(x), Inches(y), Inches(width), Inches(height)
        )

    @staticmethod
    def _write_notes(slide, page: CourseSection) -> None:
        notes = page.script
        if page.notes:
            notes = f"{notes}\n\n讲师备注：\n{page.notes}".strip()
        if notes:
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            if notes_frame is not None:
                notes_frame.text = notes
