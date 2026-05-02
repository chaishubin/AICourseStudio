"""
PPT 内容提取工具
- 提取每页文字 -> outputs/<页码>/text.txt
- 提取每页内嵌图片 -> outputs/<页码>/image_<n>.<ext>
- 每页渲染为整页截图 -> outputs/<页码>/slide.png
- 每页文字转语音 -> outputs/<页码>/audio.mp3
- 合成最终视频 -> outputs/<ppt同名>.mp4
"""

import argparse
import asyncio
import io
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image


def extract_text_from_slide(slide) -> str:
    """提取幻灯片中所有文本框的文字，保留层级缩进"""
    lines = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            level = para.level if para.level else 0
            indent = "  " * level
            lines.append(f"{indent}{text}")
    return "\n".join(lines)


def extract_images_from_slide(slide, out_dir: Path) -> list[Path]:
    """提取幻灯片中所有内嵌图片，返回保存路径列表"""
    saved = []
    img_index = 1

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.ext
            img_path = out_dir / f"image_{img_index}.{ext}"
            img_path.write_bytes(image.blob)
            saved.append(img_path)
            img_index += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = child.image
                    ext = image.ext
                    img_path = out_dir / f"image_{img_index}.{ext}"
                    img_path.write_bytes(image.blob)
                    saved.append(img_path)
                    img_index += 1

    return saved


def render_slides_to_images(ppt_path: str, output_root: Path, pages_to_render: list[int] | None = None) -> list[Path]:
    """使用 spire.presentation 将每页 PPT 渲染为 PNG 图像

    Args:
        ppt_path: PPT 文件路径
        output_root: 输出根目录
        pages_to_render: 需要渲染的页码列表（1-based），为 None 时渲染所有页
    """
    from spire.presentation import Presentation as SpirePresentation

    # 全部已存在则不加载 spire
    if pages_to_render is not None and len(pages_to_render) == 0:
        return []

    prs = SpirePresentation()
    prs.LoadFromFile(ppt_path)

    render_set = set(pages_to_render) if pages_to_render is not None else None

    slide_images = []
    for i in range(prs.Slides.Count):
        page_num = i + 1
        # 如果指定了要渲染的页，跳过不在列表中的
        if render_set is not None and page_num not in render_set:
            continue

        slide = prs.Slides[i]
        img_stream = slide.SaveAsImage()
        img_bytes = bytes(img_stream.ToArray())
        img = Image.open(io.BytesIO(img_bytes))

        page_dir = output_root / str(page_num)
        page_dir.mkdir(parents=True, exist_ok=True)
        out_path = page_dir / "slide.png"
        img.save(out_path, "PNG")
        slide_images.append(out_path)

        img_stream.Dispose()

    prs.Dispose()
    return slide_images


async def text_to_audio_async(
    texts: list[tuple[int, str, Path]],
    voice: str,
    rate: str,
) -> None:
    """并发将多页文字转换为音频文件"""
    import edge_tts

    async def convert_one(page_num: int, text: str, out_path: Path) -> None:
        content = text.strip() if text.strip() else "此页无文字内容。"
        communicate = edge_tts.Communicate(content, voice, rate=rate)
        await communicate.save(str(out_path))
        print(f"  第 {page_num} 页 音频 -> {out_path}")

    tasks = [convert_one(n, t, p) for n, t, p in texts]
    batch_size = 5
    for i in range(0, len(tasks), batch_size):
        await asyncio.gather(*tasks[i : i + batch_size])


def generate_audio(
    page_texts: list[tuple[int, str, Path]],
    voice: str = "zh-CN-XiaoxiaoNeural",
    rate: str = "+0%",
) -> None:
    asyncio.run(text_to_audio_async(page_texts, voice, rate))


def compose_video(
    output_root: Path,
    total_pages: int,
    video_path: Path,
) -> None:
    """
    将每页的 slide.png + audio.mp3 合成为完整 MP4 视频。
    每页时长 = 对应音频时长；图像静止显示。
    """
    from moviepy import ImageClip, AudioFileClip, concatenate_videoclips

    clips = []
    for i in range(1, total_pages + 1):
        page_dir = output_root / str(i)
        slide_path = page_dir / "slide.png"
        audio_path = page_dir / "audio.mp3"

        if not slide_path.exists():
            print(f"  [跳过] 第 {i} 页：缺少 slide.png", file=sys.stderr)
            continue
        if not audio_path.exists():
            print(f"  [跳过] 第 {i} 页：缺少 audio.mp3", file=sys.stderr)
            continue

        audio_clip = AudioFileClip(str(audio_path))
        duration = audio_clip.duration

        image_clip = (
            ImageClip(str(slide_path)).with_duration(duration).with_audio(audio_clip)
        )
        clips.append(image_clip)
        print(f"  第 {i} 页 片段时长: {duration:.1f}s")

    if not clips:
        print("[错误] 没有可合成的片段", file=sys.stderr)
        return

    print(f"\n合并 {len(clips)} 个片段，输出 -> {video_path}")
    final = concatenate_videoclips(clips, method="compose")
    final.write_videofile(
        str(video_path),
        fps=24,
        codec="libx264",
        audio_codec="aac",
        logger="bar",
    )
    final.close()
    for clip in clips:
        clip.close()

    size_mb = video_path.stat().st_size / 1024 / 1024
    print(f"视频生成完成：{video_path}  ({size_mb:.1f} MB)")


def process_ppt(
    ppt_path: str,
    output_root: str = "outputs",
    tts: bool = True,
    voice: str = "zh-CN-XiaoxiaoNeural",
    rate: str = "+0%",
    video: bool = True,
    skip_existing: bool = True,
) -> None:
    out_root = Path(output_root)
    out_root.mkdir(parents=True, exist_ok=True)

    prs = Presentation(ppt_path)
    total = len(prs.slides)
    print(f"共 {total} 页，开始提取...\n")

    # Step 1: 提取文字 & 内嵌图片
    page_texts: list[tuple[int, str, Path]] = []
    for i, slide in enumerate(prs.slides, start=1):
        page_dir = out_root / str(i)
        page_dir.mkdir(parents=True, exist_ok=True)

        text_path = page_dir / "text.txt"
        if skip_existing and text_path.exists():
            text = text_path.read_text(encoding="utf-8")
            print(f"  第 {i} 页 文字已存在，跳过提取  ({len(text)} 字符)")
        else:
            text = extract_text_from_slide(slide)
            text_path.write_text(text, encoding="utf-8")
            print(f"  第 {i} 页 文字 -> {text_path}  ({len(text)} 字符)")

            imgs = extract_images_from_slide(slide, page_dir)
            for p in imgs:
                print(f"  第 {i} 页 图片 -> {p}")
            if not imgs:
                print(f"  第 {i} 页 无内嵌图片")

        if tts:
            page_texts.append((i, text, page_dir / "audio.mp3"))

    # Step 2: 渲染幻灯片截图
    print("\n开始渲染幻灯片截图...")
    try:
        # 预检查哪些页缺少 slide.png
        pages_to_render = []
        for i in range(1, total + 1):
            slide_path = out_root / str(i) / "slide.png"
            if not (skip_existing and slide_path.exists()):
                pages_to_render.append(i)

        if not pages_to_render:
            print("  所有幻灯片截图已存在，跳过渲染")
        else:
            slide_paths = render_slides_to_images(ppt_path, out_root, pages_to_render)
            for i, p in enumerate(pages_to_render, start=1):
                print(f"  第 {pages_to_render[i-1]} 页 截图 -> {slide_paths[i-1]}")
            # 报告跳过的页
            skipped = total - len(pages_to_render)
            if skipped:
                print(f"  （跳过 {skipped} 页已有截图）")
    except Exception as e:
        print(f"[警告] 幻灯片渲染失败: {e}", file=sys.stderr)

    # Step 3: 文字转语音
    if tts and page_texts:
        # 过滤已有音频的页面
        if skip_existing:
            filtered = [(n, t, p) for n, t, p in page_texts if not p.exists()]
            skipped = len(page_texts) - len(filtered)
            if skipped:
                print(f"  {skipped} 页音频已存在，跳过TTS")
            page_texts = filtered

        if page_texts:
            print(f"\n开始文字转语音（声音: {voice}，语速: {rate}）...")
            try:
                generate_audio(page_texts, voice=voice, rate=rate)
            except Exception as e:
                print(f"[警告] TTS 转换失败: {e}", file=sys.stderr)
                print("  请检查网络连接，edge-tts 需要访问微软服务器", file=sys.stderr)
        else:
            print("  所有音频已存在，跳过TTS")

    # Step 4: 合成视频
    if video:
        # 输出 MP4 与 PPT 同名，放在 output_root 下
        ppt_stem = Path(ppt_path).stem
        video_path = out_root / f"{ppt_stem}.mp4"
        print(f"\n开始合成视频...")
        try:
            compose_video(out_root, total, video_path)
        except Exception as e:
            print(f"[警告] 视频合成失败: {e}", file=sys.stderr)

    print(f"\n完成！输出目录: {out_root.resolve()}")


def main():
    parser = argparse.ArgumentParser(
        description="提取 PPT 文字、图片，渲染截图，生成语音，合成视频"
    )
    parser.add_argument("ppt", help="输入 PPT/PPTX 文件路径")
    parser.add_argument(
        "-o", "--output", default="outputs", help="输出目录（默认: outputs）"
    )
    parser.add_argument("--no-tts", action="store_true", help="跳过文字转语音步骤")
    parser.add_argument("--no-video", action="store_true", help="跳过视频合成步骤")
    parser.add_argument(
        "--no-skip-existing",
        action="store_true",
        help="强制重新处理所有页面，即使输出文件已存在",
    )
    parser.add_argument(
        "--voice",
        default="zh-CN-XiaoxiaoNeural",
        help=(
            "TTS 声音角色（默认: zh-CN-XiaoxiaoNeural）\n"
            "可选中文角色:\n"
            "  zh-CN-XiaoxiaoNeural  女声·温暖（默认）\n"
            "  zh-CN-YunyangNeural   男声·专业\n"
            "  zh-CN-YunxiNeural     男声·活泼\n"
            "  zh-CN-XiaoyiNeural    女声·活泼"
        ),
    )
    parser.add_argument(
        "--rate",
        default="+0%",
        help="语速调整，如 +20%% 加快，-10%% 减慢（默认: +0%%）",
    )
    args = parser.parse_args()

    if not os.path.isfile(args.ppt):
        print(f"错误：文件不存在: {args.ppt}", file=sys.stderr)
        sys.exit(1)

    process_ppt(
        args.ppt,
        output_root=args.output,
        tts=not args.no_tts,
        voice=args.voice,
        rate=args.rate,
        video=not args.no_video,
        skip_existing=not args.no_skip_existing,
    )


if __name__ == "__main__":
    main()
